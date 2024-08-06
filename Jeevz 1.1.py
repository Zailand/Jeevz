import streamlit as st
import os
import sys
import subprocess
import pkg_resources

# Set the working directory to the script's directory
os.chdir(os.path.dirname(os.path.abspath(sys.argv[0])))

# List of required packages
required_packages = ['python-pptx', 'Pillow', 'importnb']

# Function to check and install packages
def install_packages(packages):
    for package in packages:
        try:
            dist = pkg_resources.get_distribution(package)
            print(f"{package} is already installed.")
        except pkg_resources.DistributionNotFound:
            print(f"{package} is not installed. Installing now...")
            subprocess.check_call([sys.executable, "-m", "pip", "install", package])

# Check if 'tkinter' is available (for systems that might need it)
def check_tkinter():
    try:
        import tkinter
        print("tkinter is already installed.")
    except ImportError:
        print("tkinter is not installed. Please install it through your operating system's package manager.")

# Install the required packages
install_packages(required_packages)

# Check for 'tkinter'
check_tkinter()

import os
import json
from importnb import Notebook
from pptx import Presentation
import subprocess
import pkg_resources
from tkinter import Tk, filedialog

# Function to get presentation file name from user using a file dialog
def get_presentation_filename(save=False):
    root = Tk()
    root.withdraw()  # Hide the root window
    if save:
        filename = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint files", "*.pptx")])
    else:
        filename = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
    root.destroy()
    return filename

# Function to check if a file can be saved
def can_save_file(filepath):
    try:
        # Try to save a temporary file to the same location
        temp_filepath = filepath + ".tmp"
        with open(temp_filepath, 'w') as temp_file:
            temp_file.write("test")
        os.remove(temp_filepath)
        return True
    except IOError:
        return False

# Function to load an existing presentation
def load_presentation():
    while True:
        filename = get_presentation_filename()
        if filename:
            print(f"Trying to load presentation from: {filename}")  # Debug statement
            if os.path.exists(filename):
                presentation = Presentation(filename)
                print(f"Presentation '{filename}' has been successfully loaded.")
                return presentation, filename
            else:
                print(f"File '{filename}' does not exist. Please try again.")
        else:
            print("No file selected. Please try again.")

# Function to load shared data from the notes of the first slide
def load_shared_data(presentation):
    first_slide = presentation.slides[0]
    notes_slide = first_slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text = notes_text_frame.text

    if notes_text.strip():
        return json.loads(notes_text)
    else:
        return {}

# Function to ask the user where they would like to continue from
def continue_from():
    while True:
        print("Where would you like to continue from?")
        print("1. Hypothesis, Rationale & expected results")
        print("2. Processing")
        print("3. Compression conditions")
        print("4. Tablet disintegration")
        choice = input("Enter 1, 2, 3, or 4: ")
        if choice in ['1', '2', '3', '4']:
            return int(choice)
        else:
            print("Invalid choice. Please enter 1, 2, 3, or 4.")

# Function to prompt for continuation
def continue_prompt():
    while True:
        print("\n" + "=" * 40)
        choice = input("Continue? (Y/N): ").strip().lower()
        print("=" * 40 + "\n")
        if choice in ['y', 'n']:
            return choice == 'y'
        else:
            print("Invalid choice. Please enter Y or N.")

# Import functions from Functions.ipynb
with Notebook():
    from Functions import (
        title_slide,
        hypothesis_rationale_expected_slide,
        processing_slide,
        compression_conditions_slide,
        tablet_disintegration_slide,
    )

# Function to collect user inputs and store them temporarily for an existing project
def collect_user_inputs(presentation, presentation_path, shared_data, start_from=1):
    if start_from <= 1:
        print("\n" + "=" * 40)
        print("Now working on the Hypothesis, Rationale & expected results slide")
        print("=" * 40 + "\n")
        hypothesis_rationale_expected_slide(presentation, presentation_path, shared_data)
        if not continue_prompt():
            return False

    if start_from <= 2:
        print("\n" + "=" * 40)
        print("Now working on the Processing slide")
        print("=" * 40 + "\n")
        processing_slide(presentation, presentation_path, shared_data)
        if not continue_prompt():
            return False

    if start_from <= 3:
        print("\n" + "=" * 40)
        print("Now working on the Compression conditions slide")
        print("=" * 40 + "\n")
        compression_conditions_slide(presentation, presentation_path, shared_data)
        if not continue_prompt():
            return False

    if start_from <= 4:
        print("\n" + "=" * 40)
        print("Now working on the Tablet disintegration slide")
        print("=" * 40 + "\n")
        tablet_disintegration_slide(presentation, presentation_path, shared_data)
        if not continue_prompt():
            return False

    return True

# Function to collect user inputs and store them temporarily for a new project
def collect_user_inputs_new_project(presentation, presentation_path, shared_data):
    print("\n" + "=" * 40)
    print("Now working on the Title Slide")
    print("=" * 40 + "\n")
    title_slide(presentation, presentation_path, shared_data)
    if not continue_prompt():
        return False

    if not collect_user_inputs(presentation, presentation_path, shared_data, start_from=1):
        return False

    return True

# Function to save the presentation with error handling
def save_presentation(presentation, presentation_path):
    while True:
        try:
            presentation.save(presentation_path)
            print(f"Presentation saved as {presentation_path}")
            return True
        except IOError:
            print("\n" + "*" * 80)
            print(f"*** ERROR: The file '{presentation_path}' cannot be saved. It might be open or you might not have permission. ***")
            print("*" * 80 + "\n")
            input("Please close the presentation and press Enter to try saving again...")

# Function to start a new project
def start_new_project():
    print("Starting a new project...")
    # Get the presentation filename
    presentation_path = get_presentation_filename(save=True)
    presentation = Presentation()
    shared_data = {}

    if collect_user_inputs_new_project(presentation, presentation_path, shared_data):
        while not can_save_file(presentation_path):
            print("\n" + "*" * 80)
            print(f"*** ERROR: The file '{presentation_path}' cannot be saved. It might be open or you might not have permission. ***")
            print("*" * 80 + "\n")
            input("Please close the presentation and press Enter to try saving again...")
        save_presentation(presentation, presentation_path)

# Function to load an existing project
def load_existing_project():
    print("Loading an existing project...")
    presentation, presentation_path = load_presentation()
    if not presentation:
        return
    shared_data = load_shared_data(presentation)
    start_from = continue_from()

    if collect_user_inputs(presentation, presentation_path, shared_data, start_from):
        while not can_save_file(presentation_path):
            print("\n" + "*" * 80)
            print(f"*** ERROR: The file '{presentation_path}' cannot be saved. It might be open or you might not have permission. ***")
            print("*" * 80 + "\n")
            input("Please close the presentation and press Enter to try saving again...")
        save_presentation(presentation, presentation_path)

# Main function to ask the user if they want to start a new project or load an existing one
def main():
    while True:
        try:
            choice = input("Would you like to: \n1. Start a new project\n2. Load an existing project\nEnter 1 or 2: ")
            if choice == '1':
                start_new_project()
                break
            elif choice == '2':
                load_existing_project()
                break
            else:
                print("Invalid choice. Please enter 1 or 2.")
        except Exception as e:
            print("\n" + "*" * 80)
            print(f"*** ERROR: {str(e)} ***")
            print("*" * 80 + "\n")
            input("Press Enter to try again...")

# Run the main function
main()
