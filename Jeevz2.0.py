import streamlit as st
import os
import json
from importnb import Notebook
from pptx import Presentation
import pandas as pd
import base64

# Function to load an existing presentation
def load_presentation(file):
    presentation = Presentation(file)
    st.success(f"Presentation '{file.name}' has been successfully loaded.")
    return presentation

# Function to load shared data from the notes of the first slide
def load_shared_data(presentation):
    first_slide = presentation.slides[0]
    notes_slide = first_slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text = notes_text_frame.text

    if notes_text.strip():
        return json.loads(notes_text)
    else:
        return {}

# Function to ask the user where they would like to continue from
def continue_from():
    choice = st.selectbox(
        "Where would you like to continue from?",
        ["Hypothesis, Rationale & expected results", "Processing", "Compression conditions", "Tablet disintegration"]
    )
    return ["Hypothesis, Rationale & expected results", "Processing", "Compression conditions", "Tablet disintegration"].index(choice) + 1

# Function to prompt for continuation or stopping
def continue_prompt(key):
    col1, col2 = st.columns(2)
    with col1:
        continue_clicked = st.button("Continue", key=f"continue_{key}")
    with col2:
        stop_clicked = st.button("Stop", key=f"stop_{key}")
    return continue_clicked, stop_clicked

# Import functions from Functions.ipynb
with Notebook():
    from Functions import (
        title_slide,
        hypothesis_rationale_expected_slide,
        processing_slide,
        compression_conditions_slide,
        tablet_disintegration_slide,
    )

# Function to collect user inputs and store them temporarily for an existing project
def collect_user_inputs(presentation, presentation_path, shared_data, start_from=1):
    if start_from <= 1:
        st.write("Now working on the Hypothesis, Rationale & expected results slide")
        hypothesis_rationale_expected_slide(presentation, presentation_path, shared_data)
        continue_clicked, stop_clicked = continue_prompt(key="1")
        if stop_clicked:
            st.session_state['current_step'] = 0
            return False
        if continue_clicked:
            st.session_state['current_step'] = 2

    if start_from <= 2 and st.session_state.get('current_step', 1) == 2:
        st.write("Now working on the Processing slide")
        processing_slide(presentation, presentation_path, shared_data)
        continue_clicked, stop_clicked = continue_prompt(key="2")
        if stop_clicked:
            st.session_state['current_step'] = 0
            return False
        if continue_clicked:
            st.session_state['current_step'] = 3

    if start_from <= 3 and st.session_state.get('current_step', 1) == 3:
        st.write("Now working on the Compression conditions slide")
        compression_conditions_slide(presentation, presentation_path, shared_data)
        continue_clicked, stop_clicked = continue_prompt(key="3")
        if stop_clicked:
            st.session_state['current_step'] = 0
            return False
        if continue_clicked:
            st.session_state['current_step'] = 4

    if start_from <= 4 and st.session_state.get('current_step', 1) == 4:
        st.write("Now working on the Tablet disintegration slide")
        tablet_disintegration_slide(presentation, presentation_path, shared_data)
        continue_clicked, stop_clicked = continue_prompt(key="4")
        if stop_clicked:
            st.session_state['current_step'] = 0
            return False
        if continue_clicked:
            st.session_state['current_step'] = 5

    return True

# Function to collect user inputs and store them temporarily for a new project
def collect_user_inputs_new_project(presentation, presentation_path, shared_data):
    if 'current_step' not in st.session_state:
        st.session_state['current_step'] = 0

    if st.session_state['current_step'] == 0:
        st.write("Now working on the Title Slide")
        title_slide(presentation, presentation_path, shared_data)
        continue_clicked, stop_clicked = continue_prompt(key="title")
        if stop_clicked:
            st.session_state['current_step'] = 0
            return False
        if continue_clicked:
            st.session_state['current_step'] = 1

    if st.session_state['current_step'] == 1:
        st.write("Now working on the Hypothesis, Rationale & expected results slide")
        hypothesis_rationale_expected_slide(presentation, presentation_path, shared_data)
        continue_clicked, stop_clicked = continue_prompt(key="1")
        if stop_clicked:
            st.session_state['current_step'] = 0
            return False
        if continue_clicked:
            st.session_state['current_step'] = 2

    if st.session_state['current_step'] == 2:
        st.write("Now working on the Processing slide")
        processing_slide(presentation, presentation_path, shared_data)
        continue_clicked, stop_clicked = continue_prompt(key="2")
        if stop_clicked:
            st.session_state['current_step'] = 0
            return False
        if continue_clicked:
            st.session_state['current_step'] = 3

    if st.session_state['current_step'] == 3:
        st.write("Now working on the Compression conditions slide")
        compression_conditions_slide(presentation, presentation_path, shared_data)
        continue_clicked, stop_clicked = continue_prompt(key="3")
        if stop_clicked:
            st.session_state['current_step'] = 0
            return False
        if continue_clicked:
            st.session_state['current_step'] = 4

    if st.session_state['current_step'] == 4:
        st.write("Now working on the Tablet disintegration slide")
        tablet_disintegration_slide(presentation, presentation_path, shared_data)
        continue_clicked, stop_clicked = continue_prompt(key="4")
        if stop_clicked:
            st.session_state['current_step'] = 0
            return False
        if continue_clicked:
            st.session_state['current_step'] = 5

    return True

# Function to save the presentation with error handling
def save_presentation(presentation, presentation_path):
    try:
        presentation.save(presentation_path)
        st.success(f"Presentation saved as {presentation_path}")
        return True
    except IOError:
        st.error(f"The file '{presentation_path}' cannot be saved. It might be open or you might not have permission.")
        return False

# Function to provide a download link for the presentation
def provide_download_link(presentation_path):
    with open(presentation_path, "rb") as file:
        btn = st.download_button(
            label="Download Presentation",
            data=file,
            file_name=presentation_path,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    return btn

# Function to start a new project
def start_new_project():
    st.write("Starting a new project...")
    presentation_path = st.text_input("Enter the path to save the new presentation:", "new_presentation.pptx")
    presentation = Presentation()
    shared_data = {}

    if collect_user_inputs_new_project(presentation, presentation_path, shared_data):
        if save_presentation(presentation, presentation_path):
            provide_download_link(presentation_path)
    else:
        if save_presentation(presentation, presentation_path):
            provide_download_link(presentation_path)

# Function to load an existing project
def load_existing_project():
    st.write("Loading an existing project...")
    uploaded_file = st.file_uploader("Choose a PowerPoint file", type="pptx")
    if uploaded_file is not None:
        presentation = load_presentation(uploaded_file)
        shared_data = load_shared_data(presentation)
        start_from = continue_from()

        if collect_user_inputs(presentation, uploaded_file.name, shared_data, start_from):
            if save_presentation(presentation, uploaded_file.name):
                provide_download_link(uploaded_file.name)
        else:
            if save_presentation(presentation, uploaded_file.name):
                provide_download_link(uploaded_file.name)

# Main function to ask the user if they want to start a new project or load an existing one
def main():
    st.title("PowerPoint Presentation Manager")
    choice = st.radio(
        "Would you like to:",
        ("Start a new project", "Load an existing project")
    )

    if choice == "Start a new project":
        st.session_state['current_step'] = 0
        start_new_project()
    elif choice == "Load an existing project":
        st.session_state['current_step'] = 0
        load_existing_project()

# Run the main function
if __name__ == "__main__":
    main()
