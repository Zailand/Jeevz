import streamlit as st
import os
import json
import logging
from importnb import Notebook
from pptx import Presentation
import pandas as pd

# Configure logging
logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')

# Function to load an existing presentation
def load_presentation(file):
    presentation = Presentation(file)
    st.success(f"Presentation '{file.name}' has been successfully loaded.")
    logging.debug(f"Loaded presentation: {file.name}")
    return presentation

# Function to load shared data from the notes of the first slide
def load_shared_data(presentation):
    first_slide = presentation.slides[0]
    notes_slide = first_slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text = notes_text_frame.text

    if notes_text.strip():
        return json.loads(notes_text)
    else:
        return {}

# Function to ask the user where they would like to continue from
def continue_from():
    choice = st.selectbox(
        "Where would you like to continue from?",
        ["Hypothesis, Rationale & expected results", "Processing", "Compression conditions", "Tablet disintegration"]
    )
    return ["Hypothesis, Rationale & expected results", "Processing", "Compression conditions", "Tablet disintegration"].index(choice) + 1

# Function to prompt for continuation
def continue_prompt(step):
    with st.form(key=f"form_{step}"):
        col1, col2, col3 = st.columns([1, 0.1, 1])
        with col1:
            if step == 0:
                continue_button = st.form_submit_button("Continue to Hypothesis slide")
            elif step == 1:
                continue_button = st.form_submit_button("Continue to Process slide")
            elif step == 2:
                continue_button = st.form_submit_button("Continue to Compression conditions slide")
            elif step == 3:
                continue_button = st.form_submit_button("Continue to Disintegration conditions slide")
        with col2:
            st.write("or")
        with col3:
            download_button = st.form_submit_button("Download presentation")
        return continue_button, download_button

# Function to handle the download button click
def download_presentation():
    slides_dict = st.session_state.get('slides_dict', {})
    
    # Save the dictionary to a JSON file
    dict_path = 'slides_dict.json'
    with open(dict_path, 'w') as dict_file:
        json.dump(slides_dict, dict_file)
    
    # Merge all presentation steps into one final presentation
    final_presentation_path = merge_presentations()

    # Create a zip file containing the final presentation, step presentations, and the dictionary
    import zipfile
    zip_path = 'presentation_and_dict.zip'
    with zipfile.ZipFile(zip_path, 'w') as zipf:
        step_files = sorted([f for f in os.listdir() if f.startswith("new_presentation_step_") and f.endswith(".pptx")])
        for step_file in step_files:
            zipf.write(step_file)
        zipf.write(final_presentation_path, os.path.basename(final_presentation_path))
        zipf.write(dict_path, os.path.basename(dict_path))
    
    logging.debug(f"Created zip file: {zip_path} with presentation, step files, and dictionary")

    # Provide the zip file for download
    with open(zip_path, "rb") as file:
        st.download_button(
            label="Download presentation and dictionary",
            data=file,
            file_name=zip_path,
            mime="application/zip",
            key="download_button"
        )

# Function to save the presentation with a unique filename for each step
def save_presentation(presentation, step):
    new_presentation_path = f"new_presentation_step_{step}.pptx"
    presentation.save(new_presentation_path)
    logging.debug(f"Saved presentation after step {step}: {new_presentation_path}")
    return new_presentation_path

# Function to merge all presentations into one final presentation
def merge_presentations():
    final_presentation = Presentation()
    step_files = sorted([f for f in os.listdir() if f.startswith("new_presentation_step_") and f.endswith(".pptx")])

    for step_file in step_files:
        step_presentation = Presentation(step_file)
        for slide in step_presentation.slides:
            slide_elements = slide.shapes._spTree
            new_slide_layout = final_presentation.slide_layouts[5]  # Choosing a blank layout
            new_slide = final_presentation.slides.add_slide(new_slide_layout)
            for element in slide_elements:
                new_slide.shapes._spTree.insert_element_before(element, 'p:extLst')
    
    final_presentation_path = "new_presentation_final.pptx"
    final_presentation.save(final_presentation_path)
    logging.debug(f"Final merged presentation saved as: {final_presentation_path}")
    return final_presentation_path

# Import functions from Functions.ipynb
with Notebook():
    from Functions import (
        title_slide,
        hypothesis_rationale_expected_slide,
        processing_slide,
        compression_conditions_slide,
        tablet_disintegration_slide,
    )

# Function to collect user inputs and store them temporarily for an existing project
def collect_user_inputs(presentation, presentation_path, shared_data, start_from=1):
    slides_dict = st.session_state.get('slides_dict', {})

    if start_from <= 1:
        st.write("Now working on the Hypothesis, Rationale & expected results slide")
        hypothesis_rationale_expected_slide(presentation, presentation_path, shared_data)
        slides_dict['Hypothesis, Rationale & expected results'] = 'Added'
        st.session_state.slides_dict = slides_dict
        st.write(f"Slides added: {len(slides_dict)}")
        save_presentation(presentation, 1)  # Save the presentation after adding the slide
        continue_button, download_button = continue_prompt(1)
        if download_button:
            download_presentation()
        if continue_button:
            st.session_state.current_step = 2
        else:
            return False

    if start_from <= 2:
        st.write("Now working on the Processing slide")
        processing_slide(presentation, presentation_path, shared_data)
        slides_dict['Processing'] = 'Added'
        st.session_state.slides_dict = slides_dict
        st.write(f"Slides added: {len(slides_dict)}")
        save_presentation(presentation, 2)  # Save the presentation after adding the slide
        continue_button, download_button = continue_prompt(2)
        if download_button:
            download_presentation()
        if continue_button:
            st.session_state.current_step = 3
        else:
            return False

    if start_from <= 3:
        st.write("Now working on the Compression conditions slide")
        compression_conditions_slide(presentation, presentation_path, shared_data)
        slides_dict['Compression conditions'] = 'Added'
        st.session_state.slides_dict = slides_dict
        st.write(f"Slides added: {len(slides_dict)}")
        save_presentation(presentation, 3)  # Save the presentation after adding the slide
        continue_button, download_button = continue_prompt(3)
        if download_button:
            download_presentation()
        if continue_button:
            st.session_state.current_step = 4
        else:
            return False

    if start_from <= 4:
        st.write("Now working on the Tablet disintegration slide")
        tablet_disintegration_slide(presentation, presentation_path, shared_data)
        slides_dict['Tablet disintegration'] = 'Added'
        st.session_state.slides_dict = slides_dict
        st.write(f"Slides added: {len(slides_dict)}")
        save_presentation(presentation, 4)  # Save the presentation after adding the slide
        continue_button, download_button = continue_prompt(4)
        if download_button:
            download_presentation()
        if continue_button:
            st.session_state.current_step = 5
        else:
            return False

    return True

# Function to collect user inputs and store them temporarily for a new project
def collect_user_inputs_new_project(presentation, presentation_path, shared_data):
    slides_dict = st.session_state.get('slides_dict', {})

    st.write("Now working on the Title Slide")
    title_slide(presentation, presentation_path, shared_data)
    slides_dict['Title Slide'] = 'Added'
    st.session_state.slides_dict = slides_dict
    st.write(f"Slides added: {len(slides_dict)}")
    save_presentation(presentation, 0)  # Save the presentation after adding the slide
    continue_button, download_button = continue_prompt(0)
    if download_button:
        download_presentation()
    if continue_button:
        st.session_state.current_step = 1
    else:
        return False

    if not collect_user_inputs(presentation, presentation_path, shared_data, start_from=1):
        return False

    return True

# Function to start a new project
def start_new_project():
    st.write("Starting a new project...")
    presentation_path = st.text_input("Enter the path to save the new presentation:", "new_presentation.pptx")
    presentation = Presentation()
    shared_data = {}

    if 'current_step' not in st.session_state:
        st.session_state.current_step = 0

    st.session_state.presentation_path = presentation_path

    if st.session_state.current_step == 0:
        if collect_user_inputs_new_project(presentation, presentation_path, shared_data):
            save_presentation(presentation, presentation_path)
    else:
        if collect_user_inputs(presentation, presentation_path, shared_data, start_from=st.session_state.current_step):
            save_presentation(presentation, presentation_path)

# Function to load an existing project
def load_existing_project():
    st.write("Loading an existing project...")
    uploaded_file = st.file_uploader("Choose a PowerPoint file", type="pptx")
    if uploaded_file is not None:
        presentation = load_presentation(uploaded_file)
        shared_data = load_shared_data(presentation)
        start_from = continue_from()

        if 'current_step' not in st.session_state:
            st.session_state.current_step = start_from

        st.session_state.presentation_path = uploaded_file.name

        if collect_user_inputs(presentation, uploaded_file.name, shared_data, start_from=st.session_state.current_step):
            save_presentation(presentation, uploaded_file.name)

# Main function to ask the user if they want to start a new project or load an existing one
def main():
    st.title("PowerPoint Presentation Manager")
    choice = st.radio(
        "Would you like to:",
        ("Start a new project", "Load an existing project")
    )

    if choice == "Start a new project":
        start_new_project()
    elif choice == "Load an existing project":
        load_existing_project()

# Run the main function
if __name__ == "__main__":
    main()
