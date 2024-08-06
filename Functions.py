#Cell 1
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from datetime import datetime
import json

def title_slide(presentation, presentation_path, shared_data):
    # Function to get batch number from user
    def get_batch_number():
        while True:
            batch_number = input("Enter the batch number (format XXXX-XXXX-XXXX-XX): ")
            if len(batch_number) == 17 and batch_number.count('-') == 3:
                return batch_number
            else:
                print("Invalid format. Please enter the batch number in the format XXXX-XXXX-XXXX-XX.")

    # Function to get API and excipients from user
    def get_formulation_composition():
        api_code = input("Enter API Code: ")

        while True:
            api_amount = input("Enter API amount (mg/unit): ")
            if api_amount.replace('.', '', 1).isdigit():
                api_amount = float(api_amount)
                break
            else:
                print("Invalid format. Please enter a numeric value for the API amount.")

        excipients = []
        count = 1
        while True:
            excipient_name = input(f"Enter Excipient {count}: ")
            if not excipient_name:
                break

            while True:
                excipient_amount = input(f"Enter Excipient {count} amount (mg/unit): ")
                if excipient_amount.replace('.', '', 1).isdigit():
                    excipient_amount = float(excipient_amount)
                    break
                else:
                    print("Invalid format. Please enter a numeric value for the excipient amount.")

            excipients.append((excipient_name, excipient_amount))
            count += 1
        return (api_code, api_amount), excipients

    # Function to get user initials and ELN number
    def get_initials_and_eln():
        while True:
            initials = input("Enter your initials (letters only): ")
            if initials.isalpha():
                break
            else:
                print("Invalid format. Please enter letters only.")

        while True:
            eln = input("Enter ELN number (format xxxxx-xxx): ")
            if len(eln) == 9 and eln[:5].isdigit() and eln[6:].isdigit() and eln[5] == '-':
                break
            else:
                print("Invalid format. Please enter the ELN number in the format xxxxx-xxx.")

        return initials, eln

    # Function to review and edit data
    def review_and_edit(data, prompt):
        print(f"\n{prompt}")
        for key, value in data.items():
            if key == 'excipients':
                print(f"{key}:")
                for excipient_name, excipient_amount in value:
                    print(f"  - {excipient_name}: {excipient_amount} mg/unit")
            else:
                print(f"{key}: {value}")
        while True:
            edit_choice = input("Would you like to edit any field? (Y/N): ").strip().lower()
            if edit_choice == 'y':
                field_to_edit = input("Enter the field name you want to edit: ").strip()
                if field_to_edit in data:
                    if field_to_edit == 'excipients':
                        excipients = []
                        count = 1
                        while True:
                            excipient_name = input(f"Enter Excipient {count}: ")
                            if not excipient_name:
                                break

                            while True:
                                excipient_amount = input(f"Enter Excipient {count} amount (mg/unit): ")
                                if excipient_amount.replace('.', '', 1).isdigit():
                                    excipient_amount = float(excipient_amount)
                                    break
                                else:
                                    print("Invalid format. Please enter a numeric value for the excipient amount.")

                            excipients.append((excipient_name, excipient_amount))
                            count += 1
                        data['excipients'] = excipients
                    else:
                        new_value = input(f"Enter new value for {field_to_edit}: ")
                        data[field_to_edit] = new_value
                else:
                    print("Invalid field name. Please try again.")
            elif edit_choice == 'n':
                break
            else:
                print("Invalid choice. Please enter Y or N.")
        return data

    # Add a title slide layout (usually the first layout in the template)
    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)

    # Define the title text
    title_text = "Formulation Slides"

    # Access the title placeholder and set the text
    title_placeholder = slide.shapes.title
    title_placeholder.text = title_text

    # Format the title text
    title_text_frame = title_placeholder.text_frame
    p = title_text_frame.paragraphs[0]
    p.font.name = 'Verdana (Headings)'  # Ensure this font is installed
    p.font.size = Pt(33)
    p.font.color.rgb = RGBColor(0x00, 0x19, 0x65)
    p.alignment = PP_ALIGN.LEFT

    # Define the left position for the text boxes
    left = Inches(1)

    # Get the batch number from the user
    batch_number = get_batch_number()

    # Get the formulation composition from the user
    api, excipients = get_formulation_composition()
    api_code, api_amount = api
    formulation_composition = f"Formulation composition: {api_code} ({api_amount} mg/unit)"
    for excipient_name, excipient_amount in excipients:
        formulation_composition += f", {excipient_name} ({excipient_amount} mg/unit)"

    # Get the initials and ELN number from the user
    initials, eln = get_initials_and_eln()

    # Get the current date
    current_date = datetime.now().strftime("%Y-%m-%d")

    # Store the initial values in shared_data
    shared_data['batch_number'] = batch_number
    shared_data['api_code'] = api_code
    shared_data['api_amount'] = api_amount
    shared_data['excipients'] = excipients
    shared_data['initials'] = initials
    shared_data['eln'] = eln
    shared_data['date'] = current_date

    # Review and edit step
    shared_data = review_and_edit(shared_data, "Review the entered data:")

    # Update the slide content based on the edited values
    batch_number = shared_data['batch_number']
    api_code = shared_data['api_code']
    api_amount = shared_data['api_amount']
    excipients = shared_data['excipients']
    initials = shared_data['initials']
    eln = shared_data['eln']
    current_date = shared_data['date']

    # Update the batch number textbox
    top = Inches(3)  # Adjusted position to ensure no overlap with title
    width = Inches(8)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = batch_number
    p.font.name = 'Verdana (Headings)'  # Ensure this font is installed
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)  # Red color
    p.alignment = PP_ALIGN.LEFT

    # Update the formulation composition textbox
    formulation_composition = f"Formulation composition: {api_code} ({api_amount} mg/unit)"
    for excipient_name, excipient_amount in excipients:
        formulation_composition += f", {excipient_name} ({excipient_amount} mg/unit)"

    top = Inches(4)  # Adjusted position to ensure no overlap with batch number
    width = Inches(8)
    height = Inches(2.5)  # Adjust the height as needed
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    # Set the text with wrapping
    p = tf.add_paragraph()
    p.text = formulation_composition
    p.font.name = 'Verdana (Headings)'  # Ensure this font is installed
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black color
    p.alignment = PP_ALIGN.LEFT

    # Ensure the text wraps within the text box
    tf.word_wrap = True

    # Update the final box with initials, date, and ELN number
    left = Inches(3.5)  # Center the box horizontally
    top = Inches(6)  # Align the box low on the slide
    width = Inches(3)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    # Add the initials
    p = tf.add_paragraph()
    p.text = f"Initials: {initials}"
    p.font.name = 'Verdana (Headings)'
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    p.alignment = PP_ALIGN.LEFT

    # Add the date
    p = tf.add_paragraph()
    p.text = f"Date: {current_date}"
    p.font.name = 'Verdana (Headings)'
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    p.alignment = PP_ALIGN.LEFT

    # Add the ELN number
    p = tf.add_paragraph()
    p.text = f"ELN: {eln}"
    p.font.name = 'Verdana (Headings)'
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    p.alignment = PP_ALIGN.LEFT

    # Store shared data as JSON in the notes of the first slide
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = json.dumps(shared_data)

    # Save the presentation
    presentation.save(presentation_path)
    print(f"New presentation created and saved as {presentation_path}")

#Cell 2

def hypothesis_rationale_expected_slide(presentation, presentation_path, shared_data):
    from pptx.util import Pt, Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # Add a new slide with a blank layout
    slide_layout = presentation.slide_layouts[5]  # Using the blank layout
    slide = presentation.slides.add_slide(slide_layout)

    # Remove any existing placeholders
    for shape in slide.shapes:
        if shape.is_placeholder:
            sp = shape.element
            sp.getparent().remove(sp)
    
    # Define the title text for the new slide
    title_text = "Hypothesis, Rationale & expected results"

    # Add a text box for the slide title in the upper left corner
    left = Inches(0.5)  # Adjust position as needed
    top = Inches(0.3)   # Adjust position as needed
    width = Inches(8)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame

    # Set the slide title text
    p = title_frame.add_paragraph()
    p.text = title_text

    # Format the slide title text
    p.font.name = 'Verdana'  # Ensure this font is installed
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0x00, 0x19, 0x65)
    p.alignment = PP_ALIGN.LEFT

    # Prompt user for hypotheses
    hypotheses = []
    for i in range(1, 6):
        hypothesis = input(f"Briefly (1-2 sentences) summarize hypothesis {i}: ")
        if hypothesis:
            hypotheses.append(hypothesis)
        else:
            break

    # Add a text box for the hypotheses, aligned with the title on the left and right
    left = Inches(0.5)
    top = Inches(0.9)  # 0.3 inches below the title box
    width = Inches(8)
    height = Inches(4)  # Adjust the height as needed
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_frame = content_box.text_frame

    # Add hypotheses to the text box with bullet points
    for hypothesis in hypotheses:
        p = content_frame.add_paragraph()
        p.text = hypothesis
        p.font.name = 'Verdana'  # Ensure this font is installed
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        p.alignment = PP_ALIGN.LEFT
        p.level = 0  # Bullet point level
        p.bullet = True  # Enable bullet points

    # Save the presentation with the new slide
    presentation.save(presentation_path)
    print(f"New slide added and saved in the presentation as {presentation_path}.")

#Cell 3

def processing_slide(presentation, presentation_path, shared_data):
    from pptx.util import Pt, Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import os
    import re
    from tkinter import Tk, Label, Button, filedialog
    from PIL import Image, ImageTk

    # Retrieve shared variables from shared_data
    api_code = shared_data.get('api_code', '')
    api_amount = shared_data.get('api_amount', 0)
    excipients = shared_data.get('excipients', [])
    eln_number = shared_data.get('eln', '')  # Use the ELN input from cell 1

    # Add a new slide with a blank layout
    slide_layout = presentation.slide_layouts[5]  # Using the blank layout
    slide = presentation.slides.add_slide(slide_layout)

    # Remove any existing placeholders
    for shape in slide.shapes:
        if shape.is_placeholder:
            sp = shape.element
            sp.getparent().remove(sp)

    # Define the title text for the new slide
    title_text = "Processing"

    # Add a text box for the slide title in the upper left corner
    left = Inches(0.5)  # Adjust position as needed
    top = Inches(0.3)   # Adjust position as needed
    width = Inches(8)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame

    # Set the slide title text
    p = title_frame.add_paragraph()
    p.text = title_text

    # Format the slide title text
    p.font.name = 'Verdana'  # Ensure this font is installed
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0x00, 0x19, 0x65)
    p.alignment = PP_ALIGN.LEFT

    # Add a text box for the subtitle just below the title
    subtitle_text = f"ELN: {eln_number}"
    subtitle_top = Inches(0.75)  # Adjust position as needed
    subtitle_box = slide.shapes.add_textbox(left, subtitle_top, width, height)
    subtitle_frame = subtitle_box.text_frame

    # Set the subtitle text
    p_subtitle = subtitle_frame.add_paragraph()
    p_subtitle.text = subtitle_text

    # Format the subtitle text
    p_subtitle.font.name = 'Verdana'  # Ensure this font is installed
    p_subtitle.font.size = Pt(18)
    p_subtitle.font.color.rgb = RGBColor(0x00, 0x19, 0x65)
    p_subtitle.alignment = PP_ALIGN.LEFT

    # Map of image filenames to their sizes (height by width in inches)
    image_sizes = {
        "3G HME & RC with API": (4, 5.11),
        "2G SNAC platform with API": (4, 3.56),
        "3G Roller-compaction (RC-all)": (3, 3.68),
        "3G HME SNAC - NA with API": (4, 2.88)
    }

    # Function to display images and let the user select one
    def select_image(figures_folder):
        def on_select(idx):
            nonlocal selected_image
            selected_image = idx
            root.destroy()

        root = Tk()
        root.title("Select Production Method")
        selected_image = None

        figures = [f for f in os.listdir(figures_folder) if os.path.isfile(os.path.join(figures_folder, f))]

        for idx, figure in enumerate(figures):
            img_path = os.path.join(figures_folder, figure)
            try:
                img = Image.open(img_path)
                img.thumbnail((350, 350))  # Resize to larger thumbnail
                img_tk = ImageTk.PhotoImage(img)
            except Exception as e:
                print(f"Failed to load image {img_path}: {e}")
                continue

            lbl = Label(root, image=img_tk)
            lbl.image = img_tk  # Keep a reference to avoid garbage collection
            lbl.grid(row=0, column=idx)

            btn = Button(root, text=f"Select {idx + 1}", command=lambda idx=idx: on_select(idx))
            btn.grid(row=1, column=idx)

        btn_other = Button(root, text="Other", command=lambda: on_select(len(figures)))
        btn_other.grid(row=2, column=len(figures) // 2)  # Center the "Other" button

        root.mainloop()
        return selected_image

    # Check if the "Processing Method Figures" folder exists
    figures_folder = "Processing Method Figures"
    if not os.path.exists(figures_folder):
        print(f"Error: The folder '{figures_folder}' does not exist. Please check the folder name and try again.")
        return

    # List figures in the "Processing Method Figures" folder
    figures = [f for f in os.listdir(figures_folder) if os.path.isfile(os.path.join(figures_folder, f))]

    if figures:
        selected_index = select_image(figures_folder)
        if selected_index is not None:
            if selected_index == len(figures):
                schematic_available = input("Is a production schematic available? (Yes/No): ").strip().lower()
                if schematic_available in ['yes', 'y']:
                    schematic_path = filedialog.askopenfilename(title="Select Production Schematic")
                else:
                    schematic_path = None
            else:
                schematic_path = os.path.join(figures_folder, figures[selected_index])
        else:
            schematic_available = input("Is a production schematic available? (Yes/No): ").strip().lower()
            if schematic_available in ['yes', 'y']:
                schematic_path = filedialog.askopenfilename(title="Select Production Schematic")
            else:
                schematic_path = None
    else:
        schematic_available = input("Is a production schematic available? (Yes/No): ").strip().lower()
        if schematic_available in ['yes', 'y']:
            schematic_path = filedialog.askopenfilename(title="Select Production Schematic")
        else:
            schematic_path = None

    if schematic_path:
        # Extract the filename without extension
        filename = os.path.splitext(os.path.basename(schematic_path))[0]

        # Get the size for the selected image
        if filename in image_sizes:
            height_in_inches, width_in_inches = image_sizes[filename]
        else:
            # Default size if not specified
            with Image.open(schematic_path) as img:
                img_width, img_height = img.size
            width_in_inches = img_width / 96
            height_in_inches = img_height / 96

        # Add the production schematic image to the slide at its specified size
        left = Inches(0.3)
        top = Inches(2)
        slide.shapes.add_picture(schematic_path, left, top, width=Inches(width_in_inches), height=Inches(height_in_inches))

    # Define the table position and size
    left = Inches(5.7)  # Adjust position as needed
    top = Inches(2)   # Adjust position as needed
    width = Inches(4)
    height = Inches(3)

    # Calculate the number of rows needed (1 for header + 1 for API + number of excipients + 1 for Tablet Weight)
    num_rows = 1 + 1 + len(excipients) + 1
    num_cols = 2  # Two columns: one for the name and one for the amount

    # Add the table to the slide
    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

    # Set the column widths
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(2)

    # Set the header row
    table.cell(0, 0).text = 'Component'
    table.cell(0, 1).text = 'Amount (mg/unit)'
    for cell in table.rows[0].cells:
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.name = 'Verdana'
        cell.text_frame.paragraphs[0].font.size = Pt(12)

    # Add the API to the table
    table.cell(1, 0).text = api_code
    table.cell(1, 1).text = f'{api_amount} mg/unit'
    for cell in table.rows[1].cells:
        cell.text_frame.paragraphs[0].font.name = 'Verdana'
        cell.text_frame.paragraphs[0].font.size = Pt(12)

    # Add each excipient to the table
    for i, (excipient_name, excipient_amount) in enumerate(excipients, start=2):
        table.cell(i, 0).text = excipient_name
        table.cell(i, 1).text = f'{excipient_amount} mg/unit'
        for cell in table.rows[i].cells:
            cell.text_frame.paragraphs[0].font.name = 'Verdana'
            cell.text_frame.paragraphs[0].font.size = Pt(12)

    # Calculate the total tablet weight
    total_weight = api_amount + sum(float(excipient_amount) for _, excipient_amount in excipients)

    # Add the row for the total tablet weight
    rounded_weight = round(total_weight, 2)
    table.cell(num_rows - 1, 0).text = 'Tablet Weight'
    table.cell(num_rows - 1, 1).text = f'{rounded_weight} mg'
    for cell in table.rows[num_rows - 1].cells:
        cell.text_frame.paragraphs[0].font.name = 'Verdana'
        cell.text_frame.paragraphs[0].font.size = Pt(12)

    # Save the presentation with the new slide
    presentation.save(presentation_path)
    print(f"New slide added and saved in the presentation as {presentation_path}.")

#Cell 4

# Function to review and edit data
def review_and_edit(data, prompt):
    while True:
        print(f"\n{prompt}")
        numbered_fields = list(data.items())
        for idx, (key, value) in enumerate(numbered_fields, start=1):
            if key == 'excipients':
                print(f"{idx}. {key}:")
                for excipient_name, excipient_amount in value:
                    print(f"   - {excipient_name}: {excipient_amount} mg/unit")
            else:
                print(f"{idx}. {key}: {value}")

        edit_choice = input("Would you like to edit any field? (Y/N): ").strip().lower()
        if edit_choice == 'y':
            while True:
                try:
                    field_number = int(input("Enter the number corresponding to the field you want to edit: ").strip())
                    if 1 <= field_number <= len(numbered_fields):
                        field_to_edit = numbered_fields[field_number - 1][0]
                        if field_to_edit == 'excipients':
                            excipients = []
                            count = 1
                            while True:
                                excipient_name = input(f"Enter Excipient {count}: ")
                                if not excipient_name:
                                    break

                                while True:
                                    excipient_amount = input(f"Enter Excipient {count} amount (mg/unit): ")
                                    if excipient_amount.replace('.', '', 1).isdigit():
                                        excipient_amount = float(excipient_amount)
                                        break
                                    else:
                                        print("Invalid format. Please enter a numeric value for the excipient amount.")

                                excipients.append((excipient_name, excipient_amount))
                                count += 1
                            data['excipients'] = excipients
                        else:
                            new_value = input(f"Enter new value for {field_to_edit}: ")
                            data[field_to_edit] = new_value
                        break
                    else:
                        print("Invalid number. Please enter a number corresponding to a field.")
                except ValueError:
                    print("Invalid input. Please enter a valid number.")
        elif edit_choice == 'n':
            break
        else:
            print("Invalid choice. Please enter Y or N.")
    return data

def compression_conditions_slide(presentation, presentation_path, shared_data):

    from pptx.util import Pt, Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import re
    import json

    # Retrieve shared variables from shared_data
    api_code = shared_data.get('api_code', '')
    api_amount = shared_data.get('api_amount', 0)
    excipients = shared_data.get('excipients', [])

    # Add a new slide with a blank layout
    slide_layout = presentation.slide_layouts[5]  # Using the blank layout
    slide = presentation.slides.add_slide(slide_layout)

    # Remove any existing placeholders
    for shape in slide.shapes:
        if shape.is_placeholder:
            sp = shape.element
            sp.getparent().remove(sp)

    # Define the title text for the new slide
    title_text = "Compression conditions"

    # Add a text box for the slide title in the upper left corner
    left = Inches(0.5)  # Adjust position as needed
    top = Inches(0.3)   # Adjust position as needed
    width = Inches(8)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame

    # Set the slide title text
    p = title_frame.add_paragraph()
    p.text = title_text

    # Format the slide title text
    p.font.name = 'Verdana'  # Ensure this font is installed
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0x00, 0x19, 0x65)
    p.alignment = PP_ALIGN.LEFT

    # Prompt the user for the OFR ELN number
    while True:
        eln_number = input("Please enter the OFR ELN number (format xxxxx-xxx, or press Enter to skip): ")
        if eln_number == "" or re.match(r'^\d{5}-\d{3}$', eln_number):
            break
        else:
            print("Invalid ELN number format. Please enter in the format xxxxx-xxx.")

    # Add a text box for the subtitle just below the title
    subtitle_text = f"ELN: {eln_number}"
    subtitle_top = Inches(0.75)  # Adjust position as needed
    subtitle_box = slide.shapes.add_textbox(left, subtitle_top, width, height)
    subtitle_frame = subtitle_box.text_frame

    # Set the subtitle text
    p_subtitle = subtitle_frame.add_paragraph()
    p_subtitle.text = subtitle_text

    # Format the subtitle text
    p_subtitle.font.name = 'Verdana'  # Ensure this font is installed
    p_subtitle.font.size = Pt(18)
    p_subtitle.font.color.rgb = RGBColor(0x00, 0x19, 0x65)
    p_subtitle.alignment = PP_ALIGN.LEFT

    # Function to get numerical input from the user
    def get_numerical_input(prompt):
        while True:
            value = input(prompt)
            if value == "":
                return value
            try:
                float(value)
                return value
            except ValueError:
                print("Invalid input. Please enter a numerical value or press Enter to skip.")

    # Prompt user for table inputs
    punch_width = get_numerical_input("Enter punch width (mm) (numerical value): ")
    punch_length = get_numerical_input("Enter punch length (mm) (numerical value): ")
    punch_number = get_numerical_input("Enter punch number (numerical value): ")
    cycles_used = get_numerical_input("Enter cycles used (numerical value): ")
    compression_force = get_numerical_input("Enter compression force (kN) (numerical value): ")
    tablet_height = get_numerical_input("Enter tablet average height (mm) (numerical value): ")
    tablet_weight = get_numerical_input("Enter tablet average weight (mg) (numerical value): ")
    solid_fraction = get_numerical_input("Enter solid fraction (numerical value): ")

    # Store the user inputs
    shared_data['eln_number'] = eln_number
    shared_data['punch_width'] = punch_width
    shared_data['punch_length'] = punch_length
    shared_data['punch_number'] = punch_number
    shared_data['cycles_used'] = cycles_used
    shared_data['compression_force'] = compression_force
    shared_data['tablet_height'] = tablet_height
    shared_data['tablet_weight'] = tablet_weight
    shared_data['solid_fraction'] = solid_fraction

    # Review and edit step
    shared_data = review_and_edit(shared_data, "Review the entered data:")

    # Update the slide content based on the edited values
    eln_number = shared_data['eln_number']
    punch_width = shared_data['punch_width']
    punch_length = shared_data['punch_length']
    punch_number = shared_data['punch_number']
    cycles_used = shared_data['cycles_used']
    compression_force = shared_data['compression_force']
    tablet_height = shared_data['tablet_height']
    tablet_weight = shared_data['tablet_weight']
    solid_fraction = shared_data['solid_fraction']

    # Define the table position and size
    left = Inches(0.2)  # Align to the left side
    top = Inches(2)     # Adjust position as needed
    width = Inches(4)   # Adjust width as needed
    height = Inches(3)  # Adjust height as needed

    # Calculate the number of rows needed (1 for header + 7 for parameters)
    num_rows = 8
    num_cols = 2  # Two columns: one for the parameter and one for the value

    # Add the table to the slide
    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

    # Set the column widths
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(2)

    # Set the header row
    table.cell(0, 0).text = "Parameter"
    table.cell(0, 1).text = "Value"
    for cell in table.rows[0].cells:
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.name = 'Verdana'
        cell.text_frame.paragraphs[0].font.size = Pt(12)

    # Define table content
    table.cell(1, 0).text = "Punch dimensions (WxL in mm):"
    table.cell(1, 1).text = f"{punch_width} x {punch_length}"
    table.cell(2, 0).text = "Punch number:"
    table.cell(2, 1).text = punch_number
    table.cell(3, 0).text = "Cycles used:"
    table.cell(3, 1).text = cycles_used
    table.cell(4, 0).text = "Compression force (kN):"
    table.cell(4, 1).text = compression_force
    table.cell(5, 0).text = "Tablet average height (mm):"
    table.cell(5, 1).text = tablet_height
    table.cell(6, 0).text = "Tablet average weight (mg):"
    table.cell(6, 1).text = tablet_weight
    table.cell(7, 0).text = "Solid fraction:"
    table.cell(7, 1).text = solid_fraction

    # Format the table text
    for row in table.rows:
        for cell in row.cells:
            cell.text_frame.paragraphs[0].font.name = 'Verdana'
            cell.text_frame.paragraphs[0].font.size = Pt(12)

    # Define the second table position and size
    left = Inches(4.5)  # Align to the right side
    top = Inches(2)     # Adjust position as needed
    width = Inches(4.5)  # Adjust width as needed
    height = Inches(3.5)  # Adjust height as needed

    # Calculate the number of rows needed (1 for header + 1 for API + number of excipients)
    num_rows = 1 + 1 + len(excipients)
    num_cols = 4  # Four columns: Component, Theoretical content, Determined content, CV (%)

    # Add the second table to the slide
    table2 = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

    # Set the column widths
    for col in table2.columns:
        col.width = Inches(1.35)
        
    # Set the header row
    table2.cell(0, 0).text = "Component"
    table2.cell(0, 1).text = "Theoretical\n content\n (mg/unit)"
    table2.cell(0, 2).text = "Determined\n content\n (mg/unit)"
    table2.cell(0, 3).text = "CV (%)"
    for cell in table2.rows[0].cells:
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.name = 'Verdana'
        cell.text_frame.paragraphs[0].font.size = Pt(12)

    # Add the API to the table
    table2.cell(1, 0).text = api_code
    table2.cell(1, 1).text = f'{api_amount}'
    for cell in table2.rows[1].cells:
        cell.text_frame.paragraphs[0].font.name = 'Verdana'
        cell.text_frame.paragraphs[0].font.size = Pt(12)

    # Add each excipient to the table
    for i, (excipient_name, excipient_amount) in enumerate(excipients, start=2):
        table2.cell(i, 0).text = excipient_name
        table2.cell(i, 1).text = f'{excipient_amount}'
        for cell in table2.rows[i].cells:
            cell.text_frame.paragraphs[0].font.name = 'Verdana'
            cell.text_frame.paragraphs[0].font.size = Pt(12)

    # Ask if "Determined content" and "CV" are available
    determined_content_available = input("Is Determined content (mg/unit) and CV (%) available? (Yes/No): ").strip().lower()

    if determined_content_available in ['yes', 'y']:
        # Prompt user for "Determined content" and "CV (%)" for API
        determined_content_api = get_numerical_input(f"Enter Determined content (mg/unit) for {api_code} (numerical value): ")
        cv_api = get_numerical_input(f"Enter CV (%) for {api_code} (numerical value): ")
        table2.cell(1, 2).text = determined_content_api
        table2.cell(1, 3).text = cv_api

        # Prompt user for "Determined content" and "CV (%)" for each excipient
        for i, (excipient_name, _) in enumerate(excipients, start=2):
            determined_content_excipient = get_numerical_input(f"Enter Determined content (mg/unit) for {excipient_name} (numerical value): ")
            cv_excipient = get_numerical_input(f"Enter CV (%) for {excipient_name} (numerical value): ")
            table2.cell(i, 2).text = determined_content_excipient
            table2.cell(i, 3).text = cv_excipient

        # Update shared_data with determined content and CV
        shared_data['determined_content'] = {api_code: determined_content_api}
        shared_data['cv'] = {api_code: cv_api}
        for i, (excipient_name, _) in enumerate(excipients):
            shared_data['determined_content'][excipient_name] = determined_content_excipient
            shared_data['cv'][excipient_name] = cv_excipient

    else:
        # If "Determined content" and "CV" are not available, set "N/A" for all relevant cells
        for i in range(1, num_rows):
            table2.cell(i, 2).text = "N/A"
            table2.cell(i, 3).text = "N/A"

    # Ensure consistent font size for all cells in the second table
    for row in table2.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.name = 'Verdana'

    # Save shared_data as JSON in the notes section of the slide
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.clear()
    notes_text_frame.text = json.dumps(shared_data, indent=2)

    # Save the presentation with the new slide
    presentation.save(presentation_path)
    print(f"New slide added and saved in the presentation as {presentation_path}.")

#Cell 5

def tablet_disintegration_slide(presentation, presentation_path, shared_data):

    from pptx.util import Pt, Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.text import MSO_ANCHOR
    import json
    import re
    from tkinter import Tk, filedialog

    # Function to add a disintegration slide
    def add_disintegration_slide(presentation, media, volume, svd_data, image_paths, shared_data):
        # Add a new slide with a blank layout
        slide_layout = presentation.slide_layouts[5]  # Using the blank layout
        slide = presentation.slides.add_slide(slide_layout)

        # Remove any existing placeholders
        for shape in slide.shapes:
            if shape.is_placeholder:
                sp = shape.element
                sp.getparent().remove(sp)

        # Define the title text for the new slide
        title_text = "Tablet disintegration"

        # Add a text box for the slide title in the upper left corner
        left = Inches(0.5)  # Adjust position as needed
        top = Inches(0.3)   # Adjust position as needed
        width = Inches(8)
        height = Inches(1)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame

        # Set the slide title text
        p = title_frame.add_paragraph()
        p.text = title_text

        # Format the slide title text
        p.font.name = 'Verdana'  # Ensure this font is installed
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0x00, 0x19, 0x65)
        p.alignment = PP_ALIGN.LEFT

        # Add a text box for the subtitle just below the title
        subtitle_text = f"ELN: {eln_number}"
        subtitle_top = Inches(0.75)  # Adjust position as needed
        subtitle_box = slide.shapes.add_textbox(left, subtitle_top, width, height)
        subtitle_frame = subtitle_box.text_frame

        # Set the subtitle text
        p_subtitle = subtitle_frame.add_paragraph()
        p_subtitle.text = subtitle_text

        # Format the subtitle text
        p_subtitle.font.name = 'Verdana'  # Ensure this font is installed
        p_subtitle.font.size = Pt(18)
        p_subtitle.font.color.rgb = RGBColor(0x00, 0x19, 0x65)
        p_subtitle.alignment = PP_ALIGN.LEFT
    
        # Define the table position and size
        left = Inches(0.5)  # Align to the left side
        top = Inches(1.7)  # Position below the media text box
        width = Inches(8)
        height = Inches(2.5)  # Adjust the height as needed

        # Calculate the number of columns needed (1 for Times + number of SVDs)
        num_cols = 1 + len(svd_data)
        num_rows = 5  # 5 rows: header + T10, T50, T90, T100

        # Add the table to the slide
        table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

        # Set the column widths
        for col in table.columns:
            col.width = Inches(2)

        # Set the first column title with media and volume
        cell = table.cell(0, 0)
        p = cell.text_frame.add_paragraph()
        p.text = f"Media: {media}\nVolume: {volume} mL"
        p.font.name = 'Verdana'
        p.font.size = Pt(12)
        p.font.bold = True
        cell.vertical_anchor = MSO_ANCHOR.TOP  # Set vertical alignment to top

        # Set the Times column without subscript text
        times = ["T10", "T50", "T90", "T100"]
        for i, time in enumerate(times, start=1):
            cell = table.cell(i, 0)
            p = cell.text_frame.add_paragraph()
            p.text = time
            p.font.name = 'Verdana'
            p.font.size = Pt(12)
            cell.vertical_anchor = MSO_ANCHOR.TOP  # Set vertical alignment to top

        # Add column titles for SVDs
        for col_idx in range(1, num_cols):
            cell = table.cell(0, col_idx)
            p = cell.text_frame.add_paragraph()
            p.text = f"SVD {col_idx}"
            p.font.name = 'Verdana'
            p.font.size = Pt(12)
            p.font.bold = True
            cell.vertical_anchor = MSO_ANCHOR.TOP  # Set vertical alignment to top

        # Add the user-entered SVD data to the table
        for col_idx, (t10, t50, t90, t100) in enumerate(svd_data, start=1):
            table.cell(1, col_idx).text = t10
            table.cell(2, col_idx).text = t50
            table.cell(3, col_idx).text = t90
            table.cell(4, col_idx).text = t100
            for row_idx in range(1, 5):
                table.cell(row_idx, col_idx).vertical_anchor = MSO_ANCHOR.TOP  # Set vertical alignment to top

        # Format all cells in the table
        for row in table.rows:
            for cell in row.cells:
                cell.text_frame.paragraphs[0].font.name = 'Calibri'
                cell.text_frame.paragraphs[0].font.size = Pt(12)
                cell.vertical_anchor = MSO_ANCHOR.TOP  # Set vertical alignment to top

        # Add images to the slide with titles
        image_left = Inches(0.5)  # Adjust position as needed
        image_top = Inches(4.5)   # Adjust position as needed
        image_width = Inches(2.7)   # Adjust size as needed
        image_height = Inches(2.7)  # Adjust size as needed

        for idx, image_path in enumerate(image_paths, start=1):
            if idx > 3:
                print("SVD plots per slide limited to three. Please create another slide for this media.")
                break

            # Add the image first
            image = slide.shapes.add_picture(image_path, image_left, image_top, width=image_width, height=image_height)

            # Add title above the image
            title_top = image_top - Inches(0.3)  # Adjust this value to position the title just above the image
            title_box = slide.shapes.add_textbox(image_left, title_top, image_width, Inches(0.3))
            title_frame = title_box.text_frame
            p = title_frame.add_paragraph()
            p.text = f"SVD {idx}"
            p.font.name = 'Verdana'
            p.font.size = Pt(12)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            # Bring the title to the forefront
            sp = title_box.element
            sp.getparent().append(sp)

            image_left += Inches(3.1)  # Adjust spacing between images as needed

        # Save shared_data as JSON in the notes section of the slide
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.clear()
        notes_text_frame.text = json.dumps(shared_data, indent=2)

    # Main loop to add multiple disintegration slides
    while True:
        # Prompt user for media and volume
        media = input("Enter media (or press Enter to finish): ")
        if not media:
            break
        volume = input("Enter volume (mL): ")
        eln_number = input("Please enter the SVD ELN number (format xxxxx-xxx, or press Enter to skip): ")

        # Validate the ELN number format
        while eln_number and not re.match(r'^\d{5}-\d{3}$', eln_number):
            print("Invalid ELN number format. Please enter in the format xxxxx-xxx.")
            eln_number = input("Please enter the SVD ELN number (format xxxxx-xxx, or press Enter to skip): ")

        # Initialize SVD data collection
        svd_data = []
        svd_count = 1

        while True:
            t10 = input(f"Enter T10 for SVD {svd_count} (or press Enter to finish): ")
            if not t10:
                break
            t50 = input(f"Enter T50 for SVD {svd_count}: ")
            t90 = input(f"Enter T90 for SVD {svd_count}: ")
            t100 = input(f"Enter T100 for SVD {svd_count}: ")
            svd_data.append((t10, t50, t90, t100))
            svd_count += 1

        # Check if the user wants to add a disintegration curve
        image_paths = []
        add_curve = input("Would you like to add a Disintegration curve? (Yes/No): ").strip().lower()
        if add_curve in ['yes', 'y']:
            # Prompt user for image paths using a file dialog
            root = Tk()
            root.withdraw()  # Hide the root window
            for idx in range(1, svd_count):
                image_path = filedialog.askopenfilename(title=f"Select image for SVD {idx}")
                if not image_path:
                    break
                image_paths.append(image_path)

        # Store the user inputs in shared_data
        shared_data['media'] = media
        shared_data['volume'] = volume
        shared_data['svd_data'] = [{'t10': t10, 't50': t50, 't90': t90, 't100': t100} for t10, t50, t90, t100 in svd_data]
        shared_data['image_paths'] = image_paths

        # Add a new disintegration slide for the current media
        add_disintegration_slide(presentation, media, volume, svd_data, image_paths, shared_data)

    # Save the presentation with the new slides
    presentation.save(presentation_path)
    print(f"New slides with Tablet disintegration data and images added and saved in the presentation as {presentation_path}.")
