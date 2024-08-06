import streamlit as st
import os
import json
from importnb import Notebook
from pptx import Presentation
import pandas as pd

# Function to load an existing presentation
def load_presentation(file):
    presentation = Presentation(file)
    st.success(f"Presentation '{file.name}' has been successfully loaded.")
    return presentation

# Function to load shared data from the notes of the first slide
def load_shared_data(presentation):
    first_slide = presentation.slides[0]
    notes_slide = first_slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text = notes_text_frame.text

    if notes_text.strip():
        return json.loads(notes_text)
    else:
        return {}

# Function to ask the user where they would like to continue from
def continue_from():
    choice = st.selectbox(
        "Where would you like to continue from?",
        ["Hypothesis, Rationale & expected results", "Processing", "Compression conditions", "Tablet disintegration"]
    )
    return ["Hypothesis, Rationale & expected results", "Processing", "Compression conditions", "Tablet disintegration"].index(choice) + 1

# Function to prompt for continuation
def continue_prompt():
    return st.button("Continue")

# Import functions from Functions.ipynb
with Notebook():
    from Functions import (
        title_slide,
        hypothesis_rationale_expected_slide,
        processing_slide,
        compression_conditions_slide,
        tablet_disintegration_slide,
    )

# Function to collect user inputs and store them temporarily for an existing project
def collect_user_inputs(presentation, presentation_path, shared_data, start_from=1):
    if start_from <= 1:
        st.write("Now working on the Hypothesis, Rationale & expected results slide")
        hypothesis_rationale_expected_slide(presentation, presentation_path, shared_data)
        if not continue_prompt():
            return False

    if start_from <= 2:
        st.write("Now working on the Processing slide")
        processing_slide(presentation, presentation_path, shared_data)
        if not continue_prompt():
            return False

    if start_from <= 3:
        st.write("Now working on the Compression conditions slide")
        compression_conditions_slide(presentation, presentation_path, shared_data)
        if not continue_prompt():
            return False

    if start_from <= 4:
        st.write("Now working on the Tablet disintegration slide")
        tablet_disintegration_slide(presentation, presentation_path, shared_data)
        if not continue_prompt():
            return False

    return True

# Function to collect user inputs and store them temporarily for a new project
def collect_user_inputs_new_project(presentation, presentation_path, shared_data):
    st.write("Now working on the Title Slide")
    title_slide(presentation, presentation_path, shared_data)
    if not continue_prompt():
        return False

    if not collect_user_inputs(presentation, presentation_path, shared_data, start_from=1):
        return False

    return True

# Function to save the presentation with error handling
def save_presentation(presentation, presentation_path):
    try:
        presentation.save(presentation_path)
        st.success(f"Presentation saved as {presentation_path}")
        return True
    except IOError:
        st.error(f"The file '{presentation_path}' cannot be saved. It might be open or you might not have permission.")
        return False

# Function to start a new project
def start_new_project():
    st.write("Starting a new project...")
    presentation_path = st.text_input("Enter the path to save the new presentation:", "new_presentation.pptx")
    presentation = Presentation()
    shared_data = {}

    if collect_user_inputs_new_project(presentation, presentation_path, shared_data):
        save_presentation(presentation, presentation_path)

# Function to load an existing project
def load_existing_project():
    st.write("Loading an existing project...")
    uploaded_file = st.file_uploader("Choose a PowerPoint file", type="pptx")
    if uploaded_file is not None:
        presentation = load_presentation(uploaded_file)
        shared_data = load_shared_data(presentation)
        start_from = continue_from()

        if collect_user_inputs(presentation, uploaded_file.name, shared_data, start_from):
            save_presentation(presentation, uploaded_file.name)

# Main function to ask the user if they want to start a new project or load an existing one
def main():
    st.title("PowerPoint Presentation Manager")
    choice = st.radio(
        "Would you like to:",
        ("Start a new project", "Load an existing project")
    )

    if choice == "Start a new project":
        start_new_project()
    elif choice == "Load an existing project":
        load_existing_project()

# Run the main function
if __name__ == "__main__":
    main()
